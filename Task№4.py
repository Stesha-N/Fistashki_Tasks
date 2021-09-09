#Задание 1

import pandas as pd

hm_vk = pd.read_csv(r'C:\Users\User\Desktop\Задания\kennebecproject_report_vk_posts_hmrussia_2021-06-01_2021-09-07_wc41qc.csv', sep=';')
hm_ig = pd.read_csv(r'C:\Users\User\Desktop\Задания\kennebecproject_report_ig_posts_hm_2021-06-01_2021-09-07_ayz23v.csv', sep=';')
zara_vk = pd.read_csv(r'C:\Users\User\Desktop\Задания\kennebecproject_report_vk_posts_zara_2021-06-01_2021-09-07_a4htko.csv', sep=';')
zara_ig = pd.read_csv(r'C:\Users\User\Desktop\Задания\kennebecproject_report_ig_posts_zara_2021-06-01_2021-09-07_a0o202.csv', sep=';')

hm_common = pd.concat([hm_vk, hm_ig]).fillna(0)
zara_common = pd.concat([zara_vk, zara_ig]).fillna(0)

social_actions_hm = hm_common[['likes_count','comments_count', 'reposts_count']].values.sum()
social_actions_zara = zara_common[['likes_count','comments_count', 'reposts_count']].values.sum()
accum_views_hm = hm_common['views_count'].sum()
accum_views_zara = zara_common['views_count'].sum()
followers = {'hm_v': 706429, 'hm_i': 37300000,  'zara_v': 376260 , 'zara_i': 45700000 }
ER_hm = round((hm_common['likes_count'].sum()/len(hm_common) + hm_common['comments_count'].sum()/len(hm_common))/((followers['hm_v']+followers['hm_i'])/2), 3)
ER_zara = round((zara_common['likes_count'].sum()/len(zara_common) + zara_common['comments_count'].sum()/len(zara_common))/((followers['zara_v']+followers['zara_i'])/2), 3)

indicators = pd.DataFrame([['followers', followers['hm_v'] + followers['hm_i'], followers['zara_v'] + followers['zara_i']],['social_actions', social_actions_hm, social_actions_zara], ['accum_views', accum_views_hm, accum_views_zara], ['ER', ER_hm, ER_zara]], columns=['indicators', 'H&M', 'Zara'])
indicators.to_excel('Brand_indicators.xlsx', sheet_name='Sheet_name_1')

#Задание 2

from pptx import Presentation
import pandas as pd

posts = pd.read_excel(r'C:\Users\User\Desktop\YouScan_Mentions_IQOS_05092021-05092021_b0978.xlsx', converters = {'URL': str, 'Профиль': str })
pd.set_option('display.max_columns', None)
pd.set_option('display.max_rows', None)

tone_of_mentions = posts['Тональность'].value_counts()

sourses = posts['Источник'].value_counts()

top_post_positive = posts[(posts['Тональность'] == 'Позитивная') & 
                    (posts['Тип поста'] == 'Пост') & 
                    (posts['Тип источника'] == 'Соц. сеть') & 
                    (posts['{Author} UGC'] == '{Author} UGC')].sort_values('Подписчики', ascending = False).head(1)

top_post_negative = posts[(posts['Тональность'] == 'Негативная') & 
                    (posts['Тип поста'] == 'Пост') & 
                    (posts['Тип источника'] == 'Соц. сеть') & 
                    (posts['{Author} UGC'] == '{Author} UGC')].sort_values('Подписчики', ascending = False).head(1)

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Сводные данные'

cols = 2
rows = 7
left = Inches(0.5)
top = Inches(2.0)
width = Inches(4.0)
height = Inches(0.2)

table = shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.0)

table.cell(0, 0).text = 'Источники'
table.cell(0, 1).text = 'Количество упоминаний'

for x in range(1,7):
    table.cell(x, 0).text = str(sourses.index[x-1])
    table.cell(x, 1).text = str(sourses[x-1])

chart_data = CategoryChartData()
chart_data.categories = tone_of_mentions.index
chart_data.add_series('Тональность', tone_of_mentions)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(4.3), Inches(1.5), Inches(5.5), Inches(3.5), chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

tbx = shapes.add_textbox(Inches(0.5), Inches(7), 6, 2)
text_frame = tbx.text_frame
text_frame.text = 'Общее количество упоминаний: {}'.format(len(posts))
text_frame.margin_bottom = Inches(0.01)
text_frame.margin_left = 1
text_frame.vertical_anchor = MSO_ANCHOR.TOP
text_frame.word_wrap = False
text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

prs.save('test.pptx')
